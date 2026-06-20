"""
Generate a 44-page expanded defense presentation PPTX for SpatialEx/SpatialEx+.
- Unified academic style: dark-blue primary + orange accents
- Title underlines, section accent bars, card backgrounds, highlight blocks
- Chapter divider slides with large white title cards + subtitles
- Renders LaTeX formulas via matplotlib mathtext (STIX) and inserts them as images

Output: docs/SpatialEx_Defense.pptx
Existing output is backed up with a timestamp if docs/SpatialEx_Defense_27p.pptx already exists.
"""

import os
import tempfile
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt


# Configuration
OUTPUT_PATH = "docs/SpatialEx_Defense.pptx"
BACKUP_PATH = "docs/SpatialEx_Defense_27p.pptx"
IMG_DIR = "docs/image/fig3_diagnosis"
PAPER_IMG_DIR = "docs/image/from_paper"
OTHER_IMG_DIR = "docs/image"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ============================================================================
# Unified Presentation Style
# ============================================================================

DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT_ORANGE = RGBColor(0xE8, 0x5D, 0x04)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
PALE_BLUE = RGBColor(0xEF, 0xF2, 0xFF)
DARK_GRAY = RGBColor(0x2B, 0x2B, 0x2B)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

TITLE_FONT_SIZE = Pt(32)
SUBTITLE_FONT_SIZE = Pt(24)
BODY_FONT_SIZE = Pt(18)
SMALL_FONT_SIZE = Pt(15)
CAPTION_FONT_SIZE = Pt(14)

LINE_SPACING_TIGHT = 1.05
LINE_SPACING_NORMAL = 1.15
LINE_SPACING_LOOSE = 1.25

SLIDE_MARGIN_LEFT = Inches(0.6)
SLIDE_MARGIN_RIGHT = Inches(0.6)
SLIDE_MARGIN_TOP = Inches(0.35)
SLIDE_MARGIN_BOTTOM = Inches(0.35)


def backup_existing(output_path, backup_path):
    if not os.path.exists(output_path):
        return
    # Avoid overwriting an existing backup (e.g. the original 27-page version);
    # use a timestamped name if the target already exists.
    if os.path.exists(backup_path):
        from datetime import datetime
        stem, ext = os.path.splitext(backup_path)
        timestamped = f"{stem}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{ext}"
        shutil.copy(output_path, timestamped)
        print(f"Backed up existing PPT to: {timestamped}")
    else:
        shutil.copy(output_path, backup_path)
        print(f"Backed up existing PPT to: {backup_path}")


def set_slide_size(prs):
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT


def add_title_shape(slide, text, top=Inches(0.4), height=Inches(0.9),
                    left=Inches(0.6), width=Inches(12.0),
                    font_size=32, bold=True, color=DARK_BLUE,
                    align=PP_ALIGN.LEFT,
                    add_underline=True):
    """
    Add a styled title text box with an optional dark-blue underline.
    """
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(6)

    if add_underline:
        line_y = top + height - Inches(0.12)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            line_y,
            width,
            Inches(0.035)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = DARK_BLUE
        line.line.fill.background()

    return shape


def add_section_title(slide, text, top=Inches(1.5), left=Inches(0.6),
                     width=Inches(12.0), height=Inches(0.7),
                     font_size=22, color=DARK_BLUE,
                     accent_color=ACCENT_ORANGE):
    """
    Add a section title with a small orange accent bar on the left.
    """
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left,
        top + Inches(0.12),
        Inches(0.08),
        Inches(0.45)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    shape = slide.shapes.add_textbox(
        left + Inches(0.18),
        top,
        width - Inches(0.18),
        height
    )
    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT

    return shape


def _apply_emphasis(paragraph, text, words, base_color, accent_color, font_size):
    """Build paragraph runs, highlighting whole-word occurrences in orange."""
    paragraph.clear()
    if not words:
        run = paragraph.add_run()
        run.text = text
        run.font.color.rgb = base_color
        run.font.size = Pt(font_size)
        return

    # Sort longer words first to avoid partial replacement artefacts.
    sorted_words = sorted(words, key=len, reverse=True)
    segments = [(text, False)]
    for word in sorted_words:
        new_segments = []
        for seg_text, emphasized in segments:
            if emphasized or word not in seg_text:
                new_segments.append((seg_text, emphasized))
                continue
            parts = seg_text.split(word)
            for i, part in enumerate(parts):
                new_segments.append((part, False))
                if i < len(parts) - 1:
                    new_segments.append((word, True))
        segments = new_segments

    for seg_text, emphasized in segments:
        if not seg_text:
            continue
        run = paragraph.add_run()
        run.text = seg_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = accent_color if emphasized else base_color
        run.font.bold = emphasized


def add_bullet_box(slide, bullets, top=Inches(1.6), left=Inches(0.6),
                  width=Inches(12.0), height=Inches(5.4),
                  font_size=18, color=DARK_GRAY,
                  line_spacing=LINE_SPACING_NORMAL,
                  bullet_levels=None,
                  emphasize_words=None,
                  card_background=False):
    """
    Add a structured bullet box with better paragraph spacing.
    """
    if bullet_levels is None:
        bullet_levels = [0] * len(bullets)

    if emphasize_words is None:
        emphasize_words = []

    if card_background:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left - Inches(0.08),
            top - Inches(0.08),
            width + Inches(0.16),
            height + Inches(0.16)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        _apply_emphasis(p, bullet, emphasize_words, color, ACCENT_ORANGE, font_size)
        p.space_before = Pt(6 if bullet_levels[i] == 0 else 0)
        p.space_after = Pt(8)
        p.level = bullet_levels[i]
        p.line_spacing = line_spacing

    return shape


def add_two_column_bullets(slide, left_title, left_bullets,
                           right_title, right_bullets,
                           top=Inches(1.5), height=Inches(4.8),
                           font_size=17):
    """Add two side-by-side sections with styled titles and bullet points."""
    col_width = Inches(5.9)
    gap = Inches(0.6)

    add_section_title(
        slide, left_title,
        top=top,
        left=Inches(0.6),
        width=col_width,
        font_size=22,
        color=ACCENT_ORANGE
    )

    add_bullet_box(
        slide, left_bullets,
        top=top + Inches(0.65),
        left=Inches(0.6),
        width=col_width,
        height=height,
        font_size=font_size
    )

    add_section_title(
        slide, right_title,
        top=top,
        left=Inches(0.6) + col_width + gap,
        width=col_width,
        font_size=22,
        color=ACCENT_ORANGE
    )

    add_bullet_box(
        slide, right_bullets,
        top=top + Inches(0.65),
        left=Inches(0.6) + col_width + gap,
        width=col_width,
        height=height,
        font_size=font_size
    )


def add_highlight_block(slide, text, top=Inches(5.8), left=Inches(0.6),
                        width=Inches(12.0), height=Inches(0.8),
                        font_size=16, color=DARK_BLUE,
                        bg_color=PALE_BLUE):
    """
    Add a highlighted conclusion block with pale-blue background.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(10)
    p.space_after = Pt(10)

    return shape


def add_caption(slide, text, top=Inches(6.5), left=Inches(0.8),
               width=Inches(11.8), height=Inches(0.8),
               font_size=15, color=DARK_BLUE):
    """
    Add a caption below a figure or table.
    """
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT

    return shape


def add_divider_line(slide, top=Inches(5.0), left=Inches(0.6),
                    width=Inches(12.0), color=DARK_BLUE):
    """
    Add a thin horizontal divider line inside the slide content area.
    """
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        Inches(0.025)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def render_latex(latex_str, fontsize=22, dpi=200, pad=0.15):
    """Render a LaTeX math string to a PNG file and return the path.

    Uses matplotlib's mathtext with STIX fontset, which is a Times-like
    serif font suitable for scientific publishing and matches the slide's
    Times New Roman text style without requiring a full LaTeX installation.
    """
    with plt.rc_context({
        "font.family": "serif",
        "font.serif": ["Times New Roman"],
        "mathtext.fontset": "stix",
    }):
        fig = plt.figure(figsize=(6, 1))
        fig.text(0.5, 0.5, f"${latex_str}$",
                 fontsize=fontsize, ha="center", va="center")
        plt.axis("off")
        tmp_path = tempfile.mktemp(suffix=".png")
        plt.savefig(tmp_path, dpi=dpi, bbox_inches="tight", pad_inches=pad,
                    transparent=True)
        plt.close(fig)
    return tmp_path


def add_formula(slide, latex_str, left=Inches(0.8), top=Inches(1.5),
                height=Inches(0.6), fontsize=20):
    """Render LaTeX and add as picture."""
    img_path = render_latex(latex_str, fontsize=fontsize)
    from PIL import Image
    with Image.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h
    width = height * aspect
    pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    os.remove(img_path)
    return pic


def add_image_slide(slide, img_path, title_text, max_height=Inches(5.5)):
    """
    Add an image slide with title, centered figure, and caption area.
    """
    add_title_shape(
        slide, title_text,
        top=Inches(0.3),
        height=Inches(0.8),
        font_size=28,
        bold=True,
        color=DARK_BLUE
    )

    if not os.path.exists(img_path):
        add_bullet_box(
            slide, [f"Image not found: {img_path}"],
            top=Inches(2.0),
            font_size=18
        )
        return

    from PIL import Image

    with Image.open(img_path) as img:
        img_w, img_h = img.size

    max_w = Inches(12.0)
    scale = min(max_w / Inches(img_w / 150), max_height / Inches(img_h / 150))

    new_w = Inches(img_w / 150) * scale
    new_h = Inches(img_h / 150) * scale

    left = (SLIDE_WIDTH - new_w) / 2
    top = Inches(1.3) + (max_height - new_h) / 2

    slide.shapes.add_picture(img_path, left, top, width=new_w, height=new_h)


def add_table_slide(slide, title, headers, rows, top=Inches(1.4),
                    font_size_header=16, font_size_body=15,
                    row_height_factor=0.62):
    """
    Add a styled table with dark-blue header, light-gray striped rows,
    and centered text.
    """
    if title:
        add_title_shape(
            slide, title,
            top=Inches(0.3),
            height=Inches(0.8),
            font_size=28,
            bold=True,
            color=DARK_BLUE
        )

    n_rows = len(rows) + 1
    n_cols = len(headers)

    left = Inches(0.6)
    width = Inches(12.0)
    height = Inches(row_height_factor * n_rows)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.size = Pt(font_size_header)
        p.alignment = PP_ALIGN.CENTER

    # Body
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size_body)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER

            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table


def add_section_divider(slide, main_title, subtitle=""):
    """
    Create a clean chapter divider slide with dark-blue top bar,
    large white title area, and orange subtitle.
    """
    # Top blue bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # White title card
    title_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(2.2),
        Inches(11.333),
        Inches(1.6)
    )
    title_card.fill.solid()
    title_card.fill.fore_color.rgb = WHITE
    title_card.line.fill.background()

    add_title_shape(
        slide,
        main_title,
        top=Inches(2.45),
        height=Inches(1.1),
        left=Inches(1.2),
        width=Inches(11.0),
        font_size=54,
        bold=True,
        color=DARK_BLUE,
        align=PP_ALIGN.CENTER,
        add_underline=False
    )

    if subtitle:
        add_title_shape(
            slide,
            subtitle,
            top=Inches(3.75),
            height=Inches(0.7),
            left=Inches(1.2),
            width=Inches(11.0),
            font_size=22,
            bold=False,
            color=ACCENT_ORANGE,
            align=PP_ALIGN.CENTER,
            add_underline=False
        )


def create_presentation():
    # Backup existing PPT
    backup_existing(OUTPUT_PATH, BACKUP_PATH)

    prs = Presentation()
    set_slide_size(prs)
    blank_layout = prs.slide_layouts[6]

    # ================================================================
    # 1. Title
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  SLIDE_WIDTH, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    add_title_shape(slide, "SpatialEx/SpatialEx+ 复现与改进",
                    top=Inches(1.5), height=Inches(1.0),
                    font_size=40, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_title_shape(slide, "Fig.3 Panel Diagonal Integration 上的监督信号探索",
                    top=Inches(2.6), height=Inches(0.7),
                    font_size=24, bold=False, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
    add_bullet_box(slide, [
        "论文：High-parameter spatial multi-omics through histology-anchored integration (Nature Methods 2025)",
        "答辩日期：2025 年 6 月 21 日（线上，约 40 分钟）",
        "答辩人：李熹鸣",
        "项目仓库：https://github.com/965120527lxm-maker/SpatialEvo",
        "核心发现：MLP + Strict MNN 显著优于官方 HGNN/GT + Cycle"
    ], top=Inches(3.7), left=Inches(1.8), width=Inches(9.8), height=Inches(2.5),
       font_size=18, color=BLACK)

    # ================================================================
    # 2. Outline
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "答辩提纲")
    add_bullet_box(slide, [
        "1. 背景与任务：Spatial Diagonal Integration & Fig.3",
        "2. Fig.3 任务数学化：符号、Strict 协议、评价指标",
        "3. 环境搭建与复现：Fig.2 baseline、Fig.3 跑通、修复的 6 个 bug",
        "4. 改进一：Graph Transformer 替代 HGNN",
        "5. 诊断：Branch Decomposition 与信号来源分析",
        "6. 改进二：Strict MNN 伪标签替代 Cycle",
        "7. 实验结果：Official Split、per-gene、空间可视化",
        "8. 讨论：为什么 MLP + MNN 反而最强？",
        "9. 结论、局限与展望"
    ], top=Inches(1.4), font_size=21)

    # ================================================================
    # 3. Section divider - Background
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "一. 背景与任务",
                        "先介绍我们要解决什么问题，以及这个问题为什么难")

    # ================================================================
    # 4. Background
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "背景：空间组学的高参数困境")
    add_section_title(slide, "技术困境", top=Inches(1.3))
    add_bullet_box(slide, [
        "分辨率、通量、panel 大小三者不可兼得：",
        "  • 高分辨率技术（Xenium、Merfish、CosMx）通常只能测几十到几百个基因",
        "  • 大 panel 技术（Visium）分辨率低，且仍是转录组子集"
    ], top=Inches(1.85), height=Inches(1.7), font_size=18)
    add_section_title(slide, "解决思路", top=Inches(3.45))
    add_bullet_box(slide, [
        "相邻切片测互补 panel / 互补组学，再计算整合：",
        "  • 切片 1 测 panel A，切片 2 测 panel B",
        "  • 切片 1 测转录组，切片 2 测蛋白组"
    ], top=Inches(4.0), height=Inches(1.4), font_size=18)
    add_highlight_block(
        slide,
        "关键：Spatial Diagonal Integration —— 两个切片没有共测 omics 特征，"
        "必须借助空间位置或形态先验做桥接。",
        top=Inches(5.6), height=Inches(0.7), font_size=16
    )
    img_path = os.path.join(OTHER_IMG_DIR, "fig1_bottom_three_scenarios.png")
    if os.path.exists(img_path):
        from PIL import Image
        with Image.open(img_path) as img:
            img_w, img_h = img.size
        max_w = Inches(12.0)
        max_h = Inches(1.5)
        scale = min(max_w / Inches(img_w / 150), max_h / Inches(img_h / 150))
        new_w = Inches(img_w / 150) * scale
        new_h = Inches(img_h / 150) * scale
        left = (SLIDE_WIDTH - new_w) / 2
        top = Inches(5.55)
        slide.shapes.add_picture(img_path, left, top, width=new_w, height=new_h)

    # ================================================================
    # 5. Paper Fig 1
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(OTHER_IMG_DIR, "fig1_top_framework.png")
    add_image_slide(slide, img_path, "论文原图：SpatialEx/SpatialEx+ 方法框架")
    add_caption(slide, "上半部分：SpatialEx 的 H&E foundation model、HGNN encoder、DGI 对比学习 下半部分：SpatialEx+ 的数据准备与 Cycle Module Fig.3 的 panel diagonal 是论文方法的重要应用场景", top=Inches(6.5), font_size=15)

    # ================================================================
    # 6. Two types
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "两种 Diagonal Integration")
    add_two_column_bullets(slide,
        "Panel Diagonal Integration",
        [
            "同一切片技术（如都是 Xenium）",
            "不同基因 panel（panel A vs panel B）",
            "Fig.3 任务",
            "本工作聚焦于此",
            "",
            "相邻切片各测互补基因集，",
            "目标是互相预测缺失 panel。"
        ],
        "Omics Diagonal Integration",
        [
            "不同组学类型",
            "如转录组 + 蛋白、转录组 + 表观",
            "Fig.4/5 等任务",
            "本工作尚未复现",
            "",
            "相邻切片测不同模态，",
            "目标是跨模态预测。"
        ],
        top=Inches(1.4), height=Inches(4.8), font_size=18)

    # ================================================================
    # 7. Paper Fig 3
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(PAPER_IMG_DIR, "Fig3_41592_2025_2926_Fig3_HTML.png")
    add_image_slide(slide, img_path, "论文原图：Fig.3 Panel Diagonal Integration")
    add_caption(slide, "原文使用两张相邻乳腺癌切片（Rep1 / Rep2），各测互补 panel 官方方法：SpatialEx+ 通过 Cycle 约束进行 cross-panel 预测 本工作对该设定进行严格复现并重新评估", top=Inches(6.5), font_size=15)

    # ================================================================
    # 8. Section divider - Math formalization
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "二. Fig.3 任务数学化",
                        "把任务用符号严格化，明确什么可用、什么不可用")

    # ================================================================
    # 9. Strict protocol + metrics
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "Fig.3 Strict 协议与评价指标")
    add_section_title(slide, "符号与目标", top=Inches(1.25))
    add_bullet_box(slide, [
        "Slice 1：H&E 嵌入 X1，已测 panel A YA1，缺失 panel B YB1",
        "Slice 2：H&E 嵌入 X2，已测 panel B YB2，缺失 panel A YA2",
        "目标：预测缺失 panel"
    ], top=Inches(1.75), height=Inches(1.35), font_size=17,
       emphasize_words=["Slice 1", "Slice 2", "缺失 panel"])
    add_section_title(slide, "Strict 限制", top=Inches(3.1))
    add_bullet_box(slide, [
        "训练时不可用 held-out panel（YB1、YA2）",
        "  • 这是 Fig.3 区别于普通多任务预测的关键",
        "  • 若直接用 held-out panel 监督，问题退化为有监督回归"
    ], top=Inches(3.6), height=Inches(1.4), font_size=17,
       emphasize_words=["held-out panel"])
    add_section_title(slide, "可用信号与指标", top=Inches(5.0))
    add_bullet_box(slide, [
        "信号：H&E 形态先验 / 跨切片 pseudo-label / 自监督约束",
        "指标：gene-level PCC（主指标）、SSIM、CMD"
    ], top=Inches(5.5), height=Inches(1.0), font_size=17)
    add_formula(slide, r"\text{Slice 1: }\{X_1, Y_A^1, Y_B^1\} \to \hat{Y}_B^1",
                left=Inches(7.5), top=Inches(1.25), height=Inches(0.42), fontsize=17)
    add_formula(slide, r"\text{Slice 2: }\{X_2, Y_B^2, Y_A^2\} \to \hat{Y}_A^2",
                left=Inches(7.5), top=Inches(1.82), height=Inches(0.42), fontsize=17)
    add_formula(slide, r"\text{Goal: } \hat{Y}_B^1 \approx Y_B^1, \; \hat{Y}_A^2 \approx Y_A^2",
                left=Inches(7.5), top=Inches(2.40), height=Inches(0.42), fontsize=17)
    add_formula(slide, r"\text{Constraint: } Y_B^1, Y_A^2 \notin \mathcal{D}_{\text{train}}",
                left=Inches(7.5), top=Inches(2.98), height=Inches(0.42), fontsize=17)

    # ================================================================
    # 10. fig3b
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(OTHER_IMG_DIR, "fig3b.png")
    add_image_slide(slide, img_path, "复现图：Fig.3 Strict 协议数据划分")
    add_caption(slide, "左侧 Slice1：X1 + YA1 → predict YB1 右侧 Slice2：X2 + YB2 → predict YA2 红色框表示 held-out、训练不可用的 panel", top=Inches(6.6), font_size=15)

    # ================================================================
    # 11. Section divider - Reproduction
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "三. 环境搭建与复现",
                        "先把官方代码跑通，并补充 Fig.2 baseline 对比")

    # ================================================================
    # 12. SpatialEx+ overview
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "SpatialEx / SpatialEx+ 概述")
    add_bullet_box(slide, [
        "SpatialEx：UNI + HGNN + DGI",
        "  • UNI：病理图像预训练视觉大模型，提取 H&E 嵌入",
        "  • HGNN：在细胞和伪 spot 上构建超图，做空间聚合",
        "  • DGI：自监督对比损失，增强嵌入质量",
        "  • 训练目标：从 H&E 嵌入重建 panel 表达",
        "SpatialEx+：加入 Omics Cycle Module",
        "  • 两个 panel-specific encoder：HA 预测 panel A，HB 预测 panel B",
        "  • 中间加入 regression mapping heads，实现 YA ↔ YB 循环映射",
        "Cycle 的潜在问题：",
        "  • Cycle 只保证自洽，不保证预测接近真实 missing panel",
        "  • 若 H&E → panel 映射有偏差，Cycle 会放大该偏差"
    ], top=Inches(1.3), font_size=18,
       bullet_levels=[0, 1, 1, 1, 1, 0, 1, 1, 0, 1, 1])
    add_formula(slide, r"Y_A \to \hat{Y}_B \to Y_A' \approx Y_A",
                left=Inches(7.8), top=Inches(3.1), height=Inches(0.5), fontsize=18)
    add_formula(slide, r"Y_B \to \hat{Y}_A \to Y_B' \approx Y_B",
                left=Inches(7.8), top=Inches(3.75), height=Inches(0.5), fontsize=18)

    # ================================================================
    # 13. Paper Fig 2
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(PAPER_IMG_DIR, "Fig2_41592_2025_2926_Fig2_HTML.png")
    add_image_slide(slide, img_path, "论文原图：SpatialEx/SpatialEx+ 方法框架")
    add_caption(slide, "(a) 输入：H&E 图像、空间坐标、基因表达矩阵 (b) Hypergraph 构建：细胞和 pseudo-spot 组织成超边 (c) HGNN + DGI 训练框架；(d) SpatialEx+ 的 Cycle Module", top=Inches(6.5), font_size=15)

    # ================================================================
    # 14. Bugs
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "复现中修复的 6 个关键 Bug")
    headers = ["#", "位置", "问题", "影响"]
    rows = [
        ["1", "preprocess.normalize_graph", "未定义变量 adj", "图归一化函数无法运行"],
        ["2", "preprocess", "'crs' 拼写错误", "sparse matrix 转换失败"],
        ["3", "Build_hypergraph_spatial_and_HE", "默认返回 coo", "后续期望 CSR，类型/维度错误"],
        ["4", "Model_Plus.forward", "agg_mtx 维度不匹配", "forward 矩阵乘法报错"],
        ["5", "SpatialExP.train", "Regression 维度不匹配", "Cycle module 训练崩溃"],
        ["6", "Model_Plus / Regression", "BatchNorm batch=1 崩溃", "小 batch 训练中断"]
    ]
    add_table_slide(slide, "", headers, rows, top=Inches(1.4),
                    font_size_header=15, font_size_body=14)
    add_bullet_box(slide, [
        "所有 bug 均在官方原始仓库中验证，详见 docs/report/bug_verification_original_repo.md",
        "修复这些 bug 是后续实验的前提；否则 Fig.3 主实验无法完整跑通。"
    ], top=Inches(6.2), font_size=15, color=ACCENT_ORANGE)

    # ================================================================
    # 15. Fig.2 task
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "Fig.2 复现：任务设定")
    add_bullet_box(slide, [
        "Fig.2 与 Fig.3 的区别：",
        "  • Fig.2：输入仅为 H&E，输出为全转录组（313 genes）",
        "  • 不涉及 panel split 或随机留基因",
        "  • 难度更高：模型必须 purely 从 H&E 推断表达",
        "数据：Rep1 / Rep2，各含 UNI H&E 嵌入与 Xenium RNA 表达",
        "协议：双向跨切片 train→test",
        "  • Slice A 训练 → Slice B 全部细胞评估",
        "  • 反向再做一次（B 训练 → A 评估）",
        "为什么补充 Fig.2：",
        "  • 验证 SpatialEx/GT 在纯 H&E→omics 任务上的真实能力",
        "  • 与 Fig.3 结论形成对照"
    ], top=Inches(1.3), font_size=18,
       bullet_levels=[0, 1, 1, 1, 0, 0, 1, 1, 0, 1, 1])

    # ================================================================
    # 16. Fig.2 results
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "Fig.2 复现：主结果")
    headers = ["方法", "S1 PCC", "S2 PCC", "S1 SSIM", "S2 SSIM", "S1 CMD", "S2 CMD"]
    rows = [
        ["HGNN SpatialEx (512)", "0.257", "0.273", "0.419", "0.425", "0.205", "0.207"],
        ["DeepPT", "0.268", "0.276", "0.357", "0.366", "0.271", "0.281"],
        ["GT-512 + MFP", "0.244", "0.246", "0.339", "0.329", "0.234", "0.241"],
        ["GT-128 + MFP", "0.228", "0.236", "0.309", "0.320", "0.239", "0.252"],
        ["GT-128 + DGI", "0.225", "0.235", "0.297", "0.313", "0.238", "0.248"]
    ]
    add_table_slide(slide, "", headers, rows, top=Inches(1.4),
                    font_size_header=13, font_size_body=12,
                    row_height_factor=0.55)
    add_highlight_block(
        slide,
        "结论：Fig.2 纯 H&E→omics 任务上 HGNN 仍优于 GT；Fig.3 上 GT≈HGNN "
        "不是因为网络都够强，而是因为 H&E branch 本身信息量低。",
        top=Inches(5.55), height=Inches(0.75), font_size=16
    )

    # ================================================================
    # 17. Section divider - Network improvement
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "四. 改进一: 网络方面",
                        "怀疑网络架构是瓶颈，尝试用 Graph Transformer 替代 HGNN")

    # ================================================================
    # 18. GT attempt
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "改进一：Graph Transformer 替代 HGNN")
    add_bullet_box(slide, [
        "动机：",
        "  • HGNN 超边聚合权重固定，由超图结构决定",
        "  • GT 可用稀疏自注意力学习“哪些邻居更重要”",
        "  • 理论上更能适应不同组织区域的空间异质性",
        "实现：",
        "  • Graph Transformer Encoder 替代 HGNN",
        "  • Masked Feature Prediction（MFP）替代 DGI",
        "  • MFP：随机 mask 部分 H&E 特征，让模型从邻居重建",
        "Fig.3 结果（Cycle 监督下）：",
        "  • HGNN-512 Cycle: 0.275 / 0.301",
        "  • GT-128 Cycle:   0.267 / 0.276",
        "结论：GT 与 HGNN 基本持平，网络架构不是 Fig.3 瓶颈"
    ], top=Inches(1.3), font_size=18,
       bullet_levels=[0, 1, 1, 1, 0, 1, 1, 1, 0, 1, 1, 0])

    # ================================================================
    # 19. Signal contribution
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "08_signal_contribution.png")
    add_image_slide(slide, img_path, "诊断图：信号贡献分解")
    add_caption(slide, "定量比较 H&E、measured panel、两者融合等设置的贡献 measured panel 贡献占主导；H&E 贡献接近零甚至为负 为“网络不是瓶颈，监督信号才是关键”提供直接证据", top=Inches(6.5), font_size=15)

    # ================================================================
    # 20. GT summary
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "GT 改进小结")
    add_bullet_box(slide, [
        "把 HGNN 换成 GT 没有显著提升：",
        "  • 说明空间聚合方式不是 Fig.3 任务的瓶颈",
        "  • 自适应注意力无法弥补监督信号的不足",
        "结合 Fig.2 结果：在纯 H&E 任务上 HGNN 仍优于 GT",
        "  • 说明 Fig.3 上 GT≈HGNN 不是因为网络都够强",
        "  • 而是因为 H&E 这条输入本身信息量低",
        "下一步：显式诊断 H&E 和 measured panel 各自的作用",
        "  → 引出 Branch Decomposition 实验"
    ], top=Inches(1.45), font_size=20,
       bullet_levels=[0, 1, 1, 0, 1, 1, 0, 0])

    # ================================================================
    # 21. Section divider - Diagnosis
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "五. 有效信号诊断",
                        "网络不是瓶颈，那么真正有效的信号来自哪里？")

    # ================================================================
    # 22. Branch decomposition
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "诊断：Branch Decomposition")
    headers = ["Variant", "Slice1 PCC", "Slice2 PCC"]
    rows = [
        ["H&E branch", "-0.001", "0.010"],
        ["Panel branch", "0.016", "0.248"],
        ["0.5 average", "0.011", "0.177"],
        ["Reliability-weighted", "-0.001", "0.010"]
    ]
    add_table_slide(slide, "H&E vs Panel Branch", headers, rows, top=Inches(1.4),
                    font_size_header=17, font_size_body=16)
    add_bullet_box(slide, [
        "实验设计：把输入显式拆成两条分支",
        "  • H&E branch：只用 H&E 嵌入预测缺失 panel",
        "  • Panel branch：只用同一切片 measured panel 预测缺失 panel",
        "关键发现：",
        "  • H&E branch 接近 0：H&E 不能形成稳定的 missing panel 先验",
        "  • Panel branch dominant：有效信号来自 measured → missing panel",
        "  • Late fusion 反而拉低性能：H&E 在此任务中更像噪声而非补充"
    ], top=Inches(4.6), font_size=16, color=DARK_BLUE,
       bullet_levels=[0, 1, 1, 0, 1, 1, 1])

    # ================================================================
    # 23. Branch decomposition figure
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "01_branch_decomposition.png")
    add_image_slide(slide, img_path, "Branch Decomposition 可视化")
    add_caption(slide, "左：H&E branch 散点图接近随机；右：Panel branch 有明显线性趋势 加入 H&E 后性能下降，说明 H&E 嵌入引入了与目标无关的方差", top=Inches(6.6), font_size=15)

    # ================================================================
    # 24. Signal contribution
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "08_signal_contribution.png")
    add_image_slide(slide, img_path, "信号贡献定量分析")
    add_caption(slide, "从另一角度验证 Branch Decomposition 结论 条形图/热图展示各信号来源的贡献比例 measured panel 贡献远高于 H&E；late fusion 无法提升性能", top=Inches(6.6), font_size=15)

    # ================================================================
    # 25. Per-gene scatter
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "03_per_gene_scatter.png")
    add_image_slide(slide, img_path, "诊断图：Per-Gene 预测散点")
    add_caption(slide, "横轴：基因；纵轴：PCC 或某种相关性 比较 Cycle-only、MNN、MNN+Cycle 等方法 大多数基因上 MNN 监督优于 Cycle 监督", top=Inches(6.6), font_size=15)

    # ================================================================
    # 26. Section divider - Supervision improvement
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "六. 改进二: 监督信号方面",
                        "H&E 和 Cycle 都不够强，转向跨切片 pseudo-label")

    # ================================================================
    # 27. Cycle trap
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "Cycle Consistency 的自洽陷阱")
    add_bullet_box(slide, [
        "Cycle 约束的定义：",
        "问题：Cycle 只保证自洽，不保证预测接近真实 missing panel",
        "  • 模型可以学到一种“幻觉映射”让循环闭合",
        "  • 缺少外部真实监督时，无法约束预测靠近真实生物学信号",
        "类比：翻译模型把英语→法语→英语，句子通顺≠法语翻译准确",
        "实验证据：",
        "  • MLP + Cycle only 的 PCC 仅 0.005 / 0.013，接近随机",
        "  • 说明 Cycle 本身几乎不提供有效监督"
    ], top=Inches(1.3), font_size=18,
       bullet_levels=[0, 0, 1, 1, 0, 0, 1, 1])
    add_formula(slide, r"Y_A \to \hat{Y}_B \to Y_A' \approx Y_A",
                left=Inches(7.8), top=Inches(1.35), height=Inches(0.5), fontsize=18)
    add_formula(slide, r"Y_B \to \hat{Y}_A \to Y_B' \approx Y_B",
                left=Inches(7.8), top=Inches(2.05), height=Inches(0.5), fontsize=18)
    add_formula(slide, r"\text{Self-consistency } \nRightarrow \hat{Y}_B \approx Y_B",
                left=Inches(6.8), top=Inches(2.75), height=Inches(0.5), fontsize=17)

    # ================================================================
    # 28. Cycle trap figure
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "07_cycle_trap.png")
    add_image_slide(slide, img_path, "Cycle Self-Consistency Trap")
    add_caption(slide, "Cycle 让模型在训练集上“自我闭环”，但测试时预测可能与真实 panel 无关 这是 Cycle 作为唯一监督信号的根本缺陷", top=Inches(6.7), font_size=15)

    # ================================================================
    # 29. MNN method
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "改进二：Strict MNN 伪标签")
    add_bullet_box(slide, [
        "核心思想：用跨切片 Mutual Nearest Neighbor 构造 pseudo-label",
        "严格限制：全程不使用 held-out panel（YB1、YA2）",
        "两步桥接算法：",
        "  Step 1（H&E bridge）：",
        "    • 在 H&E 嵌入空间找 Slice1 与 Slice2 的 MNN",
        "    • 把 Slice2 的 YB2 转移到 Slice1，得到伪标签 ͂YB1",
        "  Step 2（B-panel bridge）：",
        "    • 在 panel B 空间（YB2 vs ͂YB1）找 MNN",
        "    • 把 Slice1 的 YA1 转移到 Slice2，得到伪标签 ͂YA2",
        "模型训练：MLP 输入 measured panel，监督为 strict MNN pseudo-label"
    ], top=Inches(1.25), font_size=17,
       bullet_levels=[0, 0, 0, 1, 1, 1, 1, 1, 1, 0])
    add_formula(slide, r"\tilde{Y}_B^1 = \text{MNN}(X_1, X_2, Y_B^2)",
                left=Inches(7.0), top=Inches(3.5), height=Inches(0.48), fontsize=17)
    add_formula(slide, r"\tilde{Y}_A^2 = \text{MNN}(Y_B^2, \tilde{Y}_B^1, Y_A^1)",
                left=Inches(6.6), top=Inches(4.15), height=Inches(0.48), fontsize=17)
    add_formula(slide, r"\mathcal{L} = \|\hat{Y}_B^1 - \tilde{Y}_B^1\|_2^2 + \|\hat{Y}_A^2 - \tilde{Y}_A^2\|_2^2",
                left=Inches(4.8), top=Inches(4.85), height=Inches(0.5), fontsize=16)

    # ================================================================
    # 30. MNN pipeline figure
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "06_mnn_pipeline.png")
    add_image_slide(slide, img_path, "Strict MNN 伪标签流程")
    add_caption(slide, "流程：H&E MNN bridge 建立形态-空间对应 → Panel B MNN bridge 建立表达对应 → 生成伪标签监督 MLP", top=Inches(6.7), font_size=15)

    # ================================================================
    # 31. MNN sweep
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "02_mnn_sweep.png")
    add_image_slide(slide, img_path, "MNN 参数敏感性扫描")
    add_caption(slide, "横轴：邻居数 k 或 mnn_k；纵轴：gene-level PCC MNN 方法在一定范围内鲁棒，不是依赖某个特定超参数“凑”出来的结果", top=Inches(6.6), font_size=15)

    # ================================================================
    # 32. Latent alignment + Latent MNN table
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "跨切片 Latent Alignment 与 Latent MNN")
    add_bullet_box(slide, [
        "UMAP/t-SNE 展示 H&E 空间或 panel 空间的细胞分布",
        "比较 MNN 配对前后的跨切片重叠程度",
        "不同 matching 空间的效果对比："
    ], top=Inches(1.3), font_size=18, bullet_levels=[0, 0, 0])
    headers = ["Matching space", "Slice1 learned PCC", "Slice2 learned PCC"]
    rows = [
        ["raw measured panel", "0.015", "0.264"],
        ["PCA latent (50-d)", "0.007", "0.291"],
        ["CORAL aligned", "0.010", "0.228"]
    ]
    add_table_slide(slide, "", headers, rows, top=Inches(3.0),
                    font_size_header=16, font_size_body=15,
                    row_height_factor=0.55)
    add_bullet_box(slide, [
        "PCA latent + MNN 在 Slice2 上进一步提升到 0.291",
        "CORAL 线性对齐反而下降，可能过度抹平生物学相关的切片间差异",
        "Slice1 始终接近 0，说明该方向的信息桥本身较弱"
    ], top=Inches(5.4), font_size=16, color=DARK_BLUE)

    # ================================================================
    # 33. Section divider - Results
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "七. 实验结果",
                        "把 Strict MNN 接到不同模型上，全面评估效果")

    # ================================================================
    # 34. Main results
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "Official Split 主结果")
    headers = ["编码器", "监督", "S1 PCC", "S2 PCC", "S1 SSIM", "S2 SSIM"]
    rows = [
        ["HGNN-512", "Cycle", "0.275", "0.301", "0.308", "0.332"],
        ["GT-128", "Cycle", "0.267", "0.276", "0.345", "0.357"],
        ["MLP", "Cycle only", "0.005", "0.013", "0.114", "0.107"],
        ["MLP", "Strict MNN", "0.334", "0.371", "0.374", "0.398"],
        ["MLP", "MNN + Cycle", "0.315", "0.353", "0.344", "0.388"],
        ["GT-128", "Strict MNN", "0.258", "0.289", "0.359", "0.387"],
        ["HGNN-512", "Strict MNN", "0.234", "0.273", "0.072", "0.055"]
    ]
    add_table_slide(slide, "", headers, rows, top=Inches(1.4),
                    font_size_header=14, font_size_body=13,
                    row_height_factor=0.52)
    add_highlight_block(
        slide,
        "结论：MLP + Strict MNN（0.334 / 0.371）显著优于官方 HGNN/GT + Cycle；"
        "监督信号质量比模型复杂度更关键。",
        top=Inches(5.55), height=Inches(0.75), font_size=16
    )

    # ================================================================
    # 35. PCC distribution
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "10_pcc_distribution.png")
    add_image_slide(slide, img_path, "PCC 分布对比")
    add_caption(slide, "小提琴图/箱线图比较不同方法在所有基因上的 PCC 分布 MLP + Strict MNN 的分布整体右移；Cycle 方法分布接近 0 或负值", top=Inches(6.6), font_size=15)

    # ================================================================
    # 36. Gain distribution
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "04_mnn_gain_distribution.png")
    add_image_slide(slide, img_path, "Per-Gene 提升分布")
    add_caption(slide, "横轴：ΔPCC（MNN - Cycle）；纵轴：基因密度 大多数基因提升为正；78.7% 基因获得提升 少数基因下降，说明仍有改进空间", top=Inches(6.5), font_size=15)

    # ================================================================
    # 37. Top markers + table
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "Top Marker Genes 提升")
    add_bullet_box(slide, [
        "Slice 2 上提升显著的 marker 基因："
    ], top=Inches(1.3), font_size=18)
    headers = ["Gene", "raw kNN PCC", "MNN PCC", "提升"]
    rows = [
        ["CTLA4", "0.132", "0.332", "+0.200"],
        ["PTPRC", "-0.019", "0.135", "+0.154"],
        ["ESR1", "0.197", "0.323", "+0.126"],
        ["CLEC14A", "0.453", "0.590", "+0.137"]
    ]
    add_table_slide(slide, "", headers, rows, top=Inches(2.0),
                    font_size_header=16, font_size_body=15,
                    row_height_factor=0.55)
    add_bullet_box(slide, [
        "这些基因涉及免疫、上皮/基质、血管等生物学过程",
        "配图：docs/image/fig3_diagnosis/09_top_marker_gains.png",
        "MNN 的改进不只在平均指标上，也在真实空间表达结构上有视觉可辨的提升"
    ], top=Inches(5.0), font_size=16, color=DARK_BLUE)

    # ================================================================
    # 38. Marker gene figure
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(IMG_DIR, "11_marker_genes_slice2.jpg")
    add_image_slide(slide, img_path, "Slice2 Marker Gene 预测对比")
    add_caption(slide, "对比：Ground Truth vs MLP+Strict MNN vs HGNN/GT+Cycle MNN 预测的热点区域与真实值对齐更好；Cycle 预测可能过平滑或错位", top=Inches(6.6), font_size=15)

    # ================================================================
    # 39. Section divider - Discussion
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "八. 讨论",
                        "为什么一个简单 MLP 能打败复杂网络？")

    # ================================================================
    # 40. Why MLP wins
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "为什么 MLP + MNN 反而最强？")
    add_bullet_box(slide, [
        "1. 任务本质是 panel-to-panel 分子映射，不是空间聚合",
        "   同一切片上 measured panel 与 missing panel 存在生物学相关性",
        "2. MNN 已通过跨切片 matching 编码了空间/形态信息",
        "   H&E MNN 配对形态/位置相似细胞；Panel B MNN 进一步配对表达相似细胞",
        "3. HGNN/GT 的额外空间聚合会 over-smooth 或引入 H&E 噪声",
        "   Branch Decomposition 已显示 H&E 分支接近 0",
        "4. Cycle 与 MNN 冲突：自洽约束干扰可靠的外部监督",
        "   MNN 是有噪声但无偏的外部信号；Cycle 会把预测拉离 MNN 目标",
        "5. HGNN-512 + MNN 的 SSIM 极低（0.07 / 0.06），空间结构严重失真",
        "   而 MLP + MNN 的 SSIM 最高（0.37 / 0.40），空间保真最好",
        "6. 当监督信号强时，简单 MLP 足够",
        "   复杂网络需要大量有效监督才能发挥优势；在 strict 协议下 MLP 更稳健"
    ], top=Inches(1.25), font_size=16,
       bullet_levels=[0, 1, 0, 1, 0, 1, 0, 1, 0, 1, 0, 1])

    # ================================================================
    # 41. Conclusion
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "结论")
    add_bullet_box(slide, [
        "复现：跑通 SpatialEx/SpatialEx+，修复 6 个官方代码 bug",
        "复现：补充 Fig.2 baseline，HGNN 在纯 H&E 任务上仍优于 GT",
        "诊断：Branch Decomposition 揭示 H&E branch 弱、panel branch dominant",
        "改进：Strict MNN 伪标签替代 Cycle，提供可靠的外部监督",
        "发现：MLP + Strict MNN 效果最佳（PCC 0.334 / 0.371，SSIM 0.374 / 0.398）"
    ], top=Inches(1.35), height=Inches(3.4), font_size=20,
       emphasize_words=["MLP + Strict MNN"])
    add_highlight_block(
        slide,
        "核心洞察：缺失 panel 的监督信号质量比模型复杂度更重要；"
        "在 strict 协议下，简单模型 + 好信号 > 复杂模型 + 弱信号。",
        top=Inches(5.1), height=Inches(0.9), font_size=18
    )

    # ================================================================
    # 42. Limitations
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "局限与后续方向")
    add_two_column_bullets(slide,
        "局限",
        [
            "Slice1 方向仍较弱（official split 信息桥不对称）",
            "仅单一数据集验证",
            "未与 DeepPT 在 Fig.3 对比",
            "CNN_Reg、Hist2ST 等待实现",
            "GT 显存限制，Cross-Attention 退化为 MLP",
            "MNN pseudo-label 有噪声但未被建模"
        ],
        "后续方向",
        [
            "加入位置编码的多模态 matching",
            "在更多癌种和组织上验证",
            "与 DeepPT、CNN_Reg、Hist2ST 公平比较",
            "开发更轻量的 GT 变体",
            "MNN 配对不确定性建模",
            "扩展到 Omics Diagonal Integration"
        ],
        top=Inches(1.5), height=Inches(4.8), font_size=17)

    # ================================================================
    # 43. Appendix
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_shape(slide, "附录：论文其他 Fig. 与扩展方向")
    add_bullet_box(slide, [
        "可引用/展示的论文原图：",
        "  • Fig.4：Omics Diagonal Integration",
        "  • Fig.5：更多应用场景",
        "  • Fig.6：ablation 与扩展分析",
        "  • Extended Data Fig.1–10：方法细节与补充实验",
        "讲点：",
        "  • 论文 Fig.4–6 展示了 Omics Diagonal 等更复杂任务",
        "  • 我们的 Strict MNN 思想可推广到这些场景",
        "  • 未来可借鉴 Extended Data 中的更多分析"
    ], top=Inches(1.4), font_size=19,
       bullet_levels=[0, 1, 1, 1, 1, 0, 1, 1, 1])

    # ================================================================
    # 44. Thank you
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  SLIDE_WIDTH, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    add_title_shape(slide, "谢谢！",
                    top=Inches(2.3), height=Inches(1.3),
                    left=Inches(0.6), width=Inches(12.0),
                    font_size=58, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_title_shape(slide, "请各位老师斧正 / 欢迎提问",
                    top=Inches(3.9), height=Inches(0.8),
                    left=Inches(0.6), width=Inches(12.0),
                    font_size=26, bold=False, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

    # Save
    prs.save(OUTPUT_PATH)
    print(f"44-page expanded presentation saved to: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_presentation()
